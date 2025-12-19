import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io

# --- 頁面設定 ---
st.set_page_config(page_title="AI PPT Redesign Tool", layout="wide")
st.title("🎨 AI PPT 智慧版型重繪器")
st.write("上傳你的 PPT，由 AI 提供兩種不同風格的重新設計方案。")

# --- 功能函式：讀取內容 ---
def get_ppt_content(uploaded_file):
    prs = Presentation(uploaded_file)
    content = []
    for slide in prs.slides:
        slide_data = {"title": "", "body": []}
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape == slide.shapes[0]: # 假設第一個是標題
                    slide_data["title"] = shape.text
                else:
                    slide_data["body"].append(shape.text)
        content.append(slide_data)
    return content

# --- 功能函式：套用風格並產出 ---
def create_redesigned_ppt(content, style_name):
    new_prs = Presentation()
    
    # 設定風格色彩
    colors = {
        "Minimalist": RGBColor(45, 45, 45),    # 深灰
        "Tech-Future": RGBColor(0, 102, 204)   # 科技藍
    }
    bg_colors = {
        "Minimalist": RGBColor(255, 255, 255), # 白色
        "Tech-Future": RGBColor(10, 10, 25)    # 深藍黑
    }

    for data in content:
        slide_layout = new_prs.slide_layouts[1] # 使用標題+內容版面
        slide = new_prs.slides.add_slide(slide_layout)
        
        # 1. 處理標題
        title_shape = slide.shapes.title
        title_shape.text = data["title"]
        title_text_frame = title_shape.text_frame.paragraphs[0]
        title_text_frame.font.bold = True
        title_text_frame.font.color.rgb = colors.get(style_name, RGBColor(0,0,0))
        
        # 2. 處理內文
        body_shape = slide.placeholders[1]
        body_shape.text = "\n".join(data["body"])
        
        # 根據風格調整細節
        if style_name == "Tech-Future":
            # 這裡可以加入更多邏輯，例如插入科技感的裝飾線條
            pass

    ppt_io = io.BytesIO()
    new_prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# --- UI 介面 ---
uploaded_file = st.file_uploader("選擇 PPTX 檔案", type="pptx")

if uploaded_file:
    with st.spinner("正在解析內容並生成新風格..."):
        content = get_ppt_content(uploaded_file)
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.subheader("風格一：極簡商務 (Minimalist)")
            st.info("特點：高留白、黑體字、專業感十足。")
            ppt1 = create_redesigned_ppt(content, "Minimalist")
            st.download_button("下載極簡風格", data=ppt1, file_name="minimalist_design.pptx")
            
        with col2:
            st.subheader("風格二：未來科技 (Tech-Future)")
            st.success("特點：藍色調、發光元素感、適合數位轉型。")
            ppt2 = create_redesigned_ppt(content, "Tech-Future")
            st.download_button("下載科技風格", data=ppt2, file_name="tech_future_design.pptx")

    st.divider()
    st.write("🔍 **解析到的原始內容預覽：**")
    for idx, slide in enumerate(content):
        st.write(f"Slide {idx+1}: {slide['title']}")